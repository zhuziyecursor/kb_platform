from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt


OUT = Path("build/ppt/企业AI知识库平台_高管汇报优化版.pptx")
LOGO = Path("kb-portal/web/public/logo.png")


class C:
    navy = RGBColor(20, 38, 66)
    blue = RGBColor(36, 91, 150)
    teal = RGBColor(0, 150, 136)
    cyan = RGBColor(58, 185, 194)
    orange = RGBColor(242, 153, 74)
    red = RGBColor(210, 74, 74)
    green = RGBColor(69, 163, 96)
    ink = RGBColor(36, 48, 64)
    muted = RGBColor(104, 119, 137)
    light = RGBColor(246, 248, 251)
    line = RGBColor(222, 228, 236)
    white = RGBColor(255, 255, 255)
    dark_bg = RGBColor(13, 25, 44)


FONT = "PingFang SC"
FONT_LATIN = "Aptos"


def set_text(tf, text, size=18, color=C.ink, bold=False, align=None, font=FONT):
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def add_textbox(slide, x, y, w, h, text, size=18, color=C.ink, bold=False,
                align=None, valign=MSO_ANCHOR.TOP, margin=0.05):
    shape = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    shape.text_frame.margin_left = Cm(margin)
    shape.text_frame.margin_right = Cm(margin)
    shape.text_frame.margin_top = Cm(margin)
    shape.text_frame.margin_bottom = Cm(margin)
    shape.text_frame.vertical_anchor = valign
    set_text(shape.text_frame, text, size, color, bold, align)
    return shape


def add_rect(slide, x, y, w, h, fill=C.white, line=C.line, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Cm(x), Cm(y), Cm(w), Cm(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(0.7)
    return shp


def add_circle(slide, x, y, d, fill, line=None, text="", text_color=C.white, size=16):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(x), Cm(y), Cm(d), Cm(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line or fill
    shp.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(shp.text_frame, text, size=size, color=text_color, bold=True, align=PP_ALIGN.CENTER)
    return shp


def add_line(slide, x1, y1, x2, y2, color=C.line, width=1.2, arrow=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    return conn


def add_header(slide, idx, title, section="企业AI知识库平台"):
    add_textbox(slide, 1.0, 0.45, 24.5, 0.8, title, size=24, color=C.ink, bold=True)
    add_textbox(slide, 1.05, 1.25, 15, 0.45, section, size=9.5, color=C.muted)
    add_line(slide, 1.0, 1.82, 32.8, 1.82, color=C.line, width=0.8)
    add_textbox(slide, 31.8, 0.55, 1.2, 0.45, f"{idx:02d}", size=10, color=C.muted, align=PP_ALIGN.RIGHT)


def add_footer(slide, text="内部汇报 | CTO / 技术VP / 部门总经理"):
    add_textbox(slide, 1.05, 18.3, 18, 0.35, text, size=8.5, color=C.muted)


def add_bullets(slide, x, y, w, h, items, size=15, color=C.ink, bullet_color=None, gap=0.15):
    top = y
    for item in items:
        add_circle(slide, x, top + 0.08, 0.14, bullet_color or C.teal)
        tb = add_textbox(slide, x + 0.35, top, w - 0.35, 0.72, item, size=size, color=color)
        tb.text_frame.paragraphs[0].line_spacing = 1.05
        top += 0.78 + gap


def add_card(slide, x, y, w, h, title, body=None, accent=C.teal, title_size=15, body_size=11.5):
    add_rect(slide, x, y, w, h, fill=C.white, line=C.line, radius=True)
    bar = add_rect(slide, x, y, 0.12, h, fill=accent, line=accent)
    bar.line.color.rgb = accent
    add_textbox(slide, x + 0.45, y + 0.25, w - 0.7, 0.55, title, size=title_size, color=C.ink, bold=True)
    if body:
        add_textbox(slide, x + 0.45, y + 0.95, w - 0.7, h - 1.15, body, size=body_size, color=C.muted)


def add_tag(slide, x, y, text, fill=C.light, color=C.ink):
    shp = add_rect(slide, x, y, 3.2, 0.55, fill=fill, line=fill, radius=True)
    set_text(shp.text_frame, text, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)
    shp.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shp


def bullet_text(shape, items, size=13, color=C.ink):
    tf = shape.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)


def title_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = add_rect(s, 0, 0, 33.87, 19.05, fill=C.dark_bg, line=C.dark_bg)
    bg.line.color.rgb = C.dark_bg
    add_rect(s, 0, 12.4, 33.87, 6.65, fill=RGBColor(238, 244, 248), line=RGBColor(238, 244, 248))
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Cm(1.25), Cm(0.85), width=Cm(1.05), height=Cm(1.05))
    add_textbox(s, 1.3, 3.1, 21, 1.3, "企业AI知识库平台", size=36, color=C.white, bold=True)
    add_textbox(s, 1.35, 4.65, 23, 1.05, "从“能回答”到“可信问答”", size=24, color=C.cyan, bold=True)
    add_textbox(s, 1.35, 6.2, 22, 0.7, "构建面向审计行业的企业级知识沉淀、精准检索与可治理问答基础设施", size=15, color=RGBColor(211, 222, 235))
    add_tag(s, 1.35, 7.45, "RAG 平台", fill=RGBColor(30, 66, 98), color=C.white)
    add_tag(s, 4.8, 7.45, "知识治理", fill=RGBColor(30, 66, 98), color=C.white)
    add_tag(s, 8.25, 7.45, "可信溯源", fill=RGBColor(30, 66, 98), color=C.white)
    # Abstract knowledge graph
    points = [(24.0, 3.2), (28.4, 2.2), (30.5, 5.2), (25.5, 6.5), (29.2, 8.0), (22.2, 8.6)]
    for a, b in [(0, 1), (1, 2), (2, 4), (4, 3), (3, 0), (3, 5), (0, 5), (1, 3)]:
        add_line(s, *points[a], *points[b], color=RGBColor(72, 125, 163), width=1.4)
    for i, (x, y) in enumerate(points):
        add_circle(s, x, y, 0.78 if i in (0, 3) else 0.56, C.teal if i in (0, 3) else C.blue)
    add_textbox(s, 1.35, 16.8, 19, 0.5, "汇报对象：CTO / 技术VP / 部门总经理", size=11, color=C.ink)
    add_textbox(s, 25.0, 16.8, 7.6, 0.5, "2026.05", size=11, color=C.ink, align=PP_ALIGN.RIGHT)


def agenda_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, 2, "汇报主线：先讲业务必要性，再讲工程落地与决策")
    items = [
        ("01", "为什么需要企业级知识库", "通用大模型无法直接覆盖企业内部知识、权限与时效要求"),
        ("02", "现状差距在哪里", "小智/Dify 已验证需求，但检索、解析、治理能力触达天花板"),
        ("03", "KB-Platform 怎么解决", "两条主链路：文件解析-向量入库，知识检索-可信回答"),
        ("04", "价值、路线图与决策点", "用阶段化建设把平台从 MVP 推进到可规模化运营"),
    ]
    x = 2.0
    for i, (num, title, body) in enumerate(items):
        y = 3.0 + i * 3.2
        add_circle(s, x, y, 1.05, [C.teal, C.blue, C.orange, C.green][i], text=num, size=14)
        add_textbox(s, x + 1.45, y - 0.05, 10.5, 0.55, title, size=17, color=C.ink, bold=True)
        add_textbox(s, x + 1.45, y + 0.72, 21.5, 0.6, body, size=12.5, color=C.muted)
        if i < len(items) - 1:
            add_line(s, x + 0.52, y + 1.2, x + 0.52, y + 2.75, color=C.line, width=1.5)
    add_footer(s)


def llm_limits_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, 3, "通用大模型很强，但不能直接承担企业知识问答")
    add_textbox(s, 1.2, 2.4, 25.5, 0.7, "企业问答的核心要求不是“会说”，而是“答得准、可验证、守权限、可持续更新”。", size=17, color=C.ink, bold=True)
    cards = [
        ("知识更新滞后", "训练数据有截止日期，新制度、新法规、新项目经验无法自动进入模型。", C.orange),
        ("垂直语义不深", "审计、财务、合规条款往往依赖企业语境和历史材料，通用模型难以稳定理解。", C.red),
        ("内部知识不可达", "公司制度、项目复盘、操作手册、合同模板不在公开训练语料里。", C.blue),
        ("答案缺少凭据", "没有引用来源和权限校验时，业务用户无法判断答案是否可信。", C.teal),
    ]
    for i, (title, body, accent) in enumerate(cards):
        x = 1.4 + (i % 2) * 15.8
        y = 4.2 + (i // 2) * 4.8
        add_card(s, x, y, 14.3, 3.35, title, body, accent=accent, title_size=17, body_size=12.5)
    add_rect(s, 3.0, 14.6, 27.8, 1.6, fill=RGBColor(232, 244, 247), line=RGBColor(192, 224, 230), radius=True)
    add_textbox(s, 3.4, 14.95, 27, 0.7, "高管判断：企业知识库不是大模型的附属功能，而是企业 AI 应用的可信数据底座。", size=16, color=C.ink, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s)


def rag_solution_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, 4, "解决思路：给大模型配一个“企业外脑”")
    steps = [
        ("用户提问", "自然语言问题"),
        ("查询理解", "改写 / 指代消解"),
        ("知识检索", "向量 + 权限过滤"),
        ("上下文组装", "引用 + Token预算"),
        ("模型生成", "基于资料回答"),
        ("可信输出", "答案 + 来源"),
    ]
    y = 5.2
    for i, (title, sub) in enumerate(steps):
        x = 1.1 + i * 5.35
        add_circle(s, x + 1.8, y - 1.15, 0.85, [C.blue, C.teal, C.orange, C.green, C.blue, C.teal][i], text=str(i + 1), size=13)
        add_card(s, x, y, 4.45, 2.2, title, sub, accent=[C.blue, C.teal, C.orange, C.green, C.blue, C.teal][i], title_size=13.5, body_size=10.5)
        if i < len(steps) - 1:
            add_line(s, x + 4.45, y + 1.1, x + 5.15, y + 1.1, color=C.muted, width=1.3, arrow=True)
    add_textbox(s, 1.4, 9.7, 29.8, 0.7, "RAG 的本质：先从企业知识中找证据，再让模型基于证据组织答案。", size=18, color=C.ink, bold=True, align=PP_ALIGN.CENTER)
    pillars = [
        ("本地可控", "数据、模型、权限、审计可治理"),
        ("持续沉淀", "制度更新、项目经验持续进入知识库"),
        ("答案溯源", "精确到文档、页码、片段，用户可验证"),
    ]
    for i, (title, body) in enumerate(pillars):
        add_card(s, 3.0 + i * 9.7, 12.0, 8.6, 2.7, title, body, accent=[C.teal, C.orange, C.blue][i], title_size=15, body_size=11.5)
    add_footer(s)


def current_gap_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, 5, "现状判断：小智验证了价值，但还不是企业级平台")
    add_card(s, 1.2, 2.7, 9.7, 10.8, "已经验证的业务价值",
             "• 单轮 / 多轮对话\n• SSE 流式输出\n• 专项问答与智能检索\n• 会话记录、复制、收藏、点赞\n\n结论：需求真实存在，用户愿意用。", accent=C.green, title_size=16, body_size=12.5)
    add_card(s, 12.0, 2.7, 9.7, 10.8, "当前技术形态",
             "• 文件解析依赖大模型 API\n• 解析过程黑盒，不可观测\n• 文本直接入向量库\n• 单路向量召回\n• 缺少精排、权限、评测闭环\n\n结论：更接近 PoC 和需求验证器。", accent=C.orange, title_size=16, body_size=12.5)
    add_card(s, 22.8, 2.7, 9.7, 10.8, "企业级要求",
             "• 解析可控、质量可评估\n• 入库全链路可追踪\n• 多路召回与精排\n• 多租户、多维 ACL\n• 引用可验证、拒答可解释\n• Badcase 进入持续优化闭环", accent=C.blue, title_size=16, body_size=12.5)
    add_rect(s, 3.3, 15.0, 27.2, 1.2, fill=RGBColor(255, 246, 231), line=RGBColor(249, 218, 164), radius=True)
    add_textbox(s, 3.7, 15.28, 26.5, 0.5, "定位不是替代小智，而是在已验证需求的基础上，升级为“精准、可信、可治理”的企业级知识平台。", size=14.5, color=C.ink, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s)


def debt_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, 6, "四类技术债务，会直接限制规模化落地")
    data = [
        ("文件解析", "大模型 API 解析", "长文本截断、格式丢失、表格/扫描件不稳定", C.red),
        ("数据清洗", "缺少清洗环节", "页眉页脚、乱码、HTML、特殊字符进入向量库", C.orange),
        ("数据分块", "缺少可控策略", "切大噪声高，切小上下文断裂，难以稳定引用", C.blue),
        ("知识检索", "单路向量匹配", "对编号、条款、专有名词、精确数值不敏感", C.teal),
    ]
    for i, (k, now, impact, col) in enumerate(data):
        y = 3.0 + i * 3.1
        add_circle(s, 1.5, y + 0.25, 0.8, col, text=str(i + 1), size=13)
        add_textbox(s, 2.6, y, 5.0, 0.6, k, size=17, color=C.ink, bold=True)
        add_textbox(s, 8.0, y, 8.0, 0.6, now, size=13, color=C.muted)
        add_line(s, 16.2, y + 0.32, 17.2, y + 0.32, color=C.line, width=1.3, arrow=True)
        add_textbox(s, 17.6, y, 14.5, 0.9, impact, size=13.5, color=C.ink)
        add_line(s, 2.6, y + 1.35, 32.3, y + 1.35, color=C.line, width=0.6)
    add_rect(s, 2.8, 15.7, 28.1, 1.05, fill=RGBColor(232, 244, 247), line=RGBColor(192, 224, 230), radius=True)
    add_textbox(s, 3.2, 15.92, 27.4, 0.45, "高管关注点：这些不是 UI 细节，而是答案可信度、可运营性和合规可控性的基础问题。", size=13.2, color=C.ink, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s)


def architecture_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, 7, "KB-Platform 总体架构：微服务 + 异步入库 + 可信检索")
    layers = [
        ("入口与门户", ["kb-portal", "kb-gateway", "auth-adapter"], C.blue),
        ("Java 业务服务", ["ingest-service", "vector-service", "rag-service", "llm-gateway"], C.teal),
        ("Python AI 服务", ["kb-doc-processor", "rerank-service", "embedding service"], C.orange),
        ("数据与基础设施", ["PostgreSQL", "Redis", "Kafka", "Milvus", "MinIO", "Apache Tika"], C.muted),
    ]
    y = 2.8
    for idx, (title, comps, color) in enumerate(layers):
        add_textbox(s, 1.2, y + 0.35, 4.3, 0.5, title, size=13, color=C.ink, bold=True)
        x = 6.0
        for comp in comps:
            w = max(3.0, min(5.2, 1.0 + len(comp) * 0.18))
            tag = add_rect(s, x, y, w, 1.05, fill=C.white, line=color, radius=True)
            set_text(tag.text_frame, comp, size=10.5, color=color, bold=True, align=PP_ALIGN.CENTER)
            tag.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            x += w + 0.45
        if idx < len(layers) - 1:
            add_line(s, 6.2, y + 1.55, 30.5, y + 1.55, color=C.line, width=0.8)
        y += 3.15
    add_card(s, 1.5, 15.1, 9.6, 1.85, "契约优先", "OpenAPI + Kafka Schema + Milvus Collection", accent=C.blue, title_size=13.5, body_size=10.3)
    add_card(s, 12.1, 15.1, 9.6, 1.85, "异步解耦", "Kafka 削峰，上传与问答互不阻塞", accent=C.orange, title_size=13.5, body_size=10.3)
    add_card(s, 22.7, 15.1, 9.6, 1.85, "职责清晰", "入库、向量、检索、LLM 代理独立演进", accent=C.teal, title_size=13.5, body_size=10.3)
    add_footer(s)


def ingest_pipeline_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, 8, "核心链路一：文件解析到向量入库，是一条质量流水线")
    steps = [
        ("上传", "init / verify / commit"),
        ("存储", "MinIO 原文"),
        ("解析", "Tika 文本提取"),
        ("清洗", "编码 / HTML / 页眉页脚"),
        ("分块", "固定 / SMART / SMART_LLM"),
        ("向量化", "BGE 1024 维"),
        ("入库", "Milvus upsert"),
        ("READY", "状态回写"),
    ]
    y = 5.0
    for i, (a, b) in enumerate(steps):
        x = 0.85 + i * 4.05
        add_circle(s, x + 1.25, y - 1.05, 0.7, [C.blue, C.blue, C.teal, C.teal, C.orange, C.orange, C.green, C.green][i], text=str(i + 1), size=11)
        add_card(s, x, y, 3.4, 2.1, a, b, accent=[C.blue, C.blue, C.teal, C.teal, C.orange, C.orange, C.green, C.green][i], title_size=12.5, body_size=9.5)
        if i < len(steps) - 1:
            add_line(s, x + 3.4, y + 1.05, x + 3.85, y + 1.05, color=C.line, width=1.1, arrow=True)
    add_card(s, 1.3, 10.1, 10.0, 3.0, "真实实现边界",
             "ingest-service 管状态机和元数据；kb-doc-processor 负责解析、清洗、切片与 embedding；vector-service 批量写 Milvus 并回写 READY。", accent=C.teal, title_size=14, body_size=11)
    add_card(s, 12.0, 10.1, 10.0, 3.0, "质量控制点",
             "文件大小/页数限制、清洗质量分、chunk_type 推断、tags/密级/业务域/生效期等元数据随向量一起入库。", accent=C.orange, title_size=14, body_size=11)
    add_card(s, 22.7, 10.1, 10.0, 3.0, "需要持续补强",
             "OCR、表格/图表结构化、向量软下线、失败队列和更细阶段状态，是生产化治理重点。", accent=C.blue, title_size=14, body_size=11)
    add_footer(s)


def retrieval_pipeline_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, 9, "核心链路二：知识检索不是一次向量搜索，而是可信回答流水线")
    left = [
        ("1", "缓存检查"),
        ("2", "会话上下文"),
        ("3", "查询改写"),
        ("4", "查询向量化"),
        ("5", "Milvus Dense 召回"),
        ("6", "ACL / 空间过滤"),
    ]
    right = [
        ("7", "Rerank 精排"),
        ("8", "DB ACL 二次校验"),
        ("9", "Parent 回捞"),
        ("10", "拒答判断"),
        ("11", "Prompt + Token预算"),
        ("12", "LLM 生成 + 引用"),
    ]
    for col, items in enumerate([left, right]):
        x = 2.0 + col * 15.9
        for i, (num, label) in enumerate(items):
            y = 2.7 + i * 2.15
            add_circle(s, x, y, 0.65, C.teal if col == 0 else C.orange, text=num, size=10)
            add_textbox(s, x + 0.9, y - 0.05, 11.7, 0.55, label, size=14, color=C.ink, bold=True)
            if i < len(items) - 1:
                add_line(s, x + 0.32, y + 0.72, x + 0.32, y + 1.8, color=C.line, width=1)
    add_rect(s, 2.7, 16.0, 28.7, 1.0, fill=RGBColor(232, 244, 247), line=RGBColor(192, 224, 230), radius=True)
    add_textbox(s, 3.1, 16.24, 28.0, 0.45, "当前已具备：Dense 检索 + Rerank + ACL过滤 + 引用 + 流式输出 + Pipeline Trace；Phase 2 补齐 BM25/FAQ/RRF。", size=13, color=C.ink, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s)


def differentiators_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, 10, "平台级差异：从“能用”升级到“可控、可信、可运营”")
    items = [
        ("解析可控", "Tika + 清洗 + 质量评分，避免黑盒解析直接污染向量库", C.teal),
        ("分块可调", "固定分块、规则语义分块、LLM 精修，支持 Parent-Child", C.orange),
        ("检索可解释", "召回、过滤、精排、拒答、Prompt 预算全链路可追踪", C.blue),
        ("权限可治理", "租户、密级、权限组、地区、业务域、生效期进入检索链路", C.green),
        ("答案可验证", "返回引用来源，支持文档/页码/段落级追溯", C.teal),
        ("架构可演进", "微服务职责清晰，BM25、OCR、评测闭环可分阶段补齐", C.orange),
    ]
    for i, (title, body, col) in enumerate(items):
        x = 1.4 + (i % 3) * 10.7
        y = 3.0 + (i // 3) * 5.0
        add_card(s, x, y, 9.6, 3.45, title, body, accent=col, title_size=16, body_size=12)
    add_textbox(s, 2.0, 14.3, 29.8, 1.05, "这不是把开源 RAG 组件拼起来，而是把“知识生命周期、权限、质量、检索、生成、运营”做成闭环。", size=18, color=C.ink, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s)


def comparison_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, 11, "小智 / Dify / KB-Platform：定位不同，能力边界不同")
    headers = ["能力维度", "小智 / Dify 当前", "KB-Platform 当前", "下一阶段增强"]
    col_x = [1.0, 7.1, 15.6, 24.0]
    col_w = [5.6, 8.0, 8.0, 8.4]
    y0 = 2.5
    for x, w, h in zip(col_x, col_w, headers):
        add_rect(s, x, y0, w, 0.85, fill=C.navy, line=C.navy)
        add_textbox(s, x + 0.1, y0 + 0.16, w - 0.2, 0.35, h, size=10.5, color=C.white, bold=True, align=PP_ALIGN.CENTER)
    rows = [
        ("文件解析", "大模型解析，黑盒", "Tika + 清洗 + 状态机", "OCR / 表格 / 图表"),
        ("分块策略", "弱控制", "固定 + SMART + Parent-Child", "领域模板 + 自动调参"),
        ("检索链路", "单路向量", "Dense + Rerank + ACL", "BM25 + FAQ + RRF"),
        ("权限治理", "基础权限", "密级 / 权限组 / 生效期", "OBO / 用户上下文打通"),
        ("可观测性", "弱", "Pipeline Trace + Prompt预算", "SLO / 告警 / 成本面板"),
        ("评测闭环", "点赞未闭环", "Trace 与引用基础", "Badcase → 评测集 → 优化任务"),
    ]
    y = y0 + 1.0
    for r, row in enumerate(rows):
        fill = C.white if r % 2 == 0 else C.light
        for c, val in enumerate(row):
            add_rect(s, col_x[c], y, col_w[c], 1.55, fill=fill, line=C.line)
            add_textbox(s, col_x[c] + 0.18, y + 0.18, col_w[c] - 0.36, 0.9, val, size=9.8 if c else 10.5, color=C.ink if c != 1 else C.muted, bold=(c == 0))
        y += 1.55
    add_footer(s)


def risk_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, 12, "当前工程判断：MVP 已跑通，生产化还需补齐四个硬点")
    risks = [
        ("权限上下文", "当前存在 DEV 用户、租户、权限组常量", "接入 OIDC/OBO/JWT claims，形成真实权限上下文", C.red),
        ("召回能力", "当前以 Dense 为主，精确条款/编号仍有漏召风险", "引入 BM25、FAQ、条款 Fast Path、RRF 融合", C.orange),
        ("文档治理", "OCR、多模态、软下线和细阶段状态未完全闭环", "补 OCR/表格解析、向量删除、DLQ 和阶段状态表", C.blue),
        ("评测闭环", "没有体系化检索/生成/权限回归评测", "Badcase 入库，建立 Recall、忠实度、引用准确率指标", C.teal),
    ]
    for i, (title, now, next_, col) in enumerate(risks):
        x = 1.2 + (i % 2) * 15.9
        y = 3.0 + (i // 2) * 5.2
        add_card(s, x, y, 14.2, 3.8, title, f"当前：{now}\n\n动作：{next_}", accent=col, title_size=16, body_size=11.3)
    add_textbox(s, 2.3, 14.3, 29.0, 1.2, "建议原则：MVP 可以小范围试点，但不能带着开发态权限和无评测闭环进入大规模生产。", size=17, color=C.ink, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s)


def business_value_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, 13, "业务价值：把分散知识变成可复用、可治理、可审计的生产力")
    vals = [
        ("效率", "制度、流程、项目经验快速定位，减少重复咨询和人工查文档。", "目标：常见问题秒级响应", C.teal),
        ("质量", "回答基于知识库引用，减少凭经验口径不一致。", "目标：关键答案可验证", C.blue),
        ("合规", "权限、密级、有效期进入检索链路，避免越权和过期制度召回。", "目标：知识访问可审计", C.orange),
        ("沉淀", "项目复盘、专家经验和制度更新持续入库，形成组织资产。", "目标：系统越用越准", C.green),
    ]
    for i, (title, body, metric, col) in enumerate(vals):
        x = 1.4 + i * 8.05
        add_circle(s, x + 2.8, 3.0, 1.2, col, text=title[0], size=20)
        add_card(s, x, 4.6, 7.2, 4.2, title, body, accent=col, title_size=18, body_size=12)
        add_rect(s, x, 9.2, 7.2, 1.25, fill=C.light, line=C.line, radius=True)
        add_textbox(s, x + 0.25, 9.48, 6.7, 0.45, metric, size=11.5, color=C.ink, bold=True, align=PP_ALIGN.CENTER)
    add_rect(s, 2.2, 13.2, 29.5, 1.65, fill=RGBColor(232, 244, 247), line=RGBColor(192, 224, 230), radius=True)
    add_textbox(s, 2.8, 13.57, 28.4, 0.65, "平台价值不只在“问答”，而在于建立企业知识资产的生产、消费、治理和优化闭环。", size=17, color=C.ink, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s)


def roadmap_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, 14, "三阶段路线图：先可用，再好用，最后可规模化运营")
    phases = [
        ("Phase 1\nMVP 骨架", "已具备", "上传入库、解析清洗、切片、Dense 检索、Rerank、引用、SSE、Trace、Prompt预算", C.green),
        ("Phase 2\n精准检索", "下一步重点", "BM25、FAQ、RRF、多路召回、LLM Query Rewrite、OCR/表格解析、引用校验", C.orange),
        ("Phase 3\n生产治理", "规模化", "统一认证、权限上下文、评测体系、反馈闭环、成本审计、监控告警、容量治理", C.blue),
    ]
    y = 4.3
    for i, (phase, tag, body, col) in enumerate(phases):
        x = 2.0 + i * 10.4
        add_rect(s, x, y, 8.8, 7.8, fill=C.white, line=col, radius=True)
        add_circle(s, x + 3.75, y - 0.65, 1.3, col, text=str(i + 1), size=20)
        add_textbox(s, x + 0.35, y + 1.0, 8.1, 1.0, phase, size=18, color=C.ink, bold=True, align=PP_ALIGN.CENTER)
        add_tag(s, x + 2.8, y + 2.55, tag, fill=C.light, color=col)
        add_textbox(s, x + 0.65, y + 3.55, 7.5, 2.7, body, size=11.2, color=C.muted, align=PP_ALIGN.CENTER)
        if i < 2:
            add_line(s, x + 8.8, y + 3.9, x + 10.0, y + 3.9, color=C.line, width=1.5, arrow=True)
    add_footer(s)


def metrics_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, 15, "建议验收指标：用数据证明平台是否真的变好")
    groups = [
        ("入库质量", ["文档入库成功率", "5分钟内 READY 比例", "解析质量分分布", "失败重试成功率"], C.teal),
        ("检索质量", ["Recall@K", "引用准确率", "Rerank 后相关性", "无权限过滤正确率"], C.orange),
        ("回答质量", ["答案忠实度", "拒答准确率", "用户满意度", "Badcase 关闭周期"], C.blue),
        ("系统体验", ["首 Token P95", "端到端延迟 P95", "缓存命中率", "Token 成本/会话"], C.green),
    ]
    for i, (title, items, col) in enumerate(groups):
        x = 1.3 + (i % 2) * 16.0
        y = 3.0 + (i // 2) * 5.3
        add_card(s, x, y, 14.2, 3.8, title, "\n".join(f"• {it}" for it in items), accent=col, title_size=16, body_size=12)
    add_textbox(s, 2.1, 14.4, 29.7, 0.95, "没有评测体系，RAG 优化只能靠感觉；有评测体系，才能把 Badcase 变成可管理的工程任务。", size=16, color=C.ink, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s)


def resource_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, 16, "资源与协同：平台落地需要技术、业务、数据三方共建")
    lanes = [
        ("技术团队", "平台架构、服务治理、检索策略、评测框架、性能与安全", C.blue),
        ("业务专家", "知识范围、标准问答、引用校验、Badcase 评审、验收口径", C.orange),
        ("数据/运维", "文档治理、权限目录、部署监控、容量成本、SLO 管控", C.teal),
    ]
    for i, (title, body, col) in enumerate(lanes):
        add_card(s, 2.0 + i * 10.4, 3.2, 8.8, 4.5, title, body, accent=col, title_size=17, body_size=12)
    add_rect(s, 3.0, 9.7, 27.8, 4.3, fill=C.light, line=C.line, radius=True)
    add_textbox(s, 3.55, 10.2, 27.0, 0.5, "试点建议", size=17, color=C.ink, bold=True)
    add_bullets(s, 3.65, 11.0, 25.8, 2.2, [
        "选择 1-2 个高价值知识域：制度合规、审计方法、项目复盘或操作手册。",
        "每个知识域准备标准问题集、权威答案、可引用来源，作为验收基线。",
        "试点期间按周复盘 Badcase，形成检索、分块、Prompt、权限四类优化任务。",
    ], size=12.5, bullet_color=C.teal)
    add_footer(s)


def decision_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, 17, "需要管理层明确的三个决策")
    decisions = [
        ("试点范围", "先聚焦一个业务域和一批高质量文档，避免一开始全量铺开。", C.teal),
        ("安全边界", "在真实权限上下文打通前，仅用于受控试点和非敏感/授权数据。", C.red),
        ("路线投入", "Phase 2 重点投向多路召回、OCR/表格解析、评测闭环和权限打通。", C.orange),
    ]
    for i, (title, body, col) in enumerate(decisions):
        y = 3.4 + i * 3.9
        add_circle(s, 2.0, y, 1.0, col, text=str(i + 1), size=15)
        add_card(s, 3.4, y - 0.3, 27.5, 2.55, title, body, accent=col, title_size=17, body_size=12.5)
    add_textbox(s, 3.0, 16.0, 27.5, 0.7, "建议结论：批准进入受控业务试点，并以 Phase 2 能力补齐作为生产化前置条件。", size=18, color=C.ink, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s)


def closing_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, 33.87, 19.05, fill=C.dark_bg, line=C.dark_bg)
    add_textbox(s, 2.2, 3.6, 28, 1.0, "从“能回答”到“可信问答”", size=34, color=C.white, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, 4.2, 5.25, 24, 0.8, "企业 AI 的核心不是让模型更会聊天，而是让企业知识能够被安全、准确、持续地使用。", size=17, color=RGBColor(211, 222, 235), align=PP_ALIGN.CENTER)
    summary = [
        "一条可运行的知识入库流水线",
        "一条可追踪的 RAG 检索生成流水线",
        "一个可继续演进的企业知识平台底座",
    ]
    for i, text in enumerate(summary):
        add_circle(s, 6.0, 8.2 + i * 1.6, 0.55, [C.teal, C.orange, C.blue][i], text="✓", size=13)
        add_textbox(s, 6.9, 8.15 + i * 1.6, 20.8, 0.5, text, size=16, color=C.white, bold=True)
    add_textbox(s, 2.2, 16.2, 29.5, 0.6, "谢谢", size=20, color=C.cyan, bold=True, align=PP_ALIGN.CENTER)


def build():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)
    for fn in [
        title_slide,
        agenda_slide,
        llm_limits_slide,
        rag_solution_slide,
        current_gap_slide,
        debt_slide,
        architecture_slide,
        ingest_pipeline_slide,
        retrieval_pipeline_slide,
        differentiators_slide,
        comparison_slide,
        risk_slide,
        business_value_slide,
        roadmap_slide,
        metrics_slide,
        resource_slide,
        decision_slide,
        closing_slide,
    ]:
        fn(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
